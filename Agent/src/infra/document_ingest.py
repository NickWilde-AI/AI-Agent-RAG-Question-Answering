"""统一文档入库：PDF / Office / 表格 -> 页面级 Page。"""

from __future__ import annotations

import csv
import shutil
import subprocess
import tempfile
import warnings
from pathlib import Path
from typing import Callable, Iterable, List, Optional

from ..models import Page
from .pdf_ingest import ingest_pdf_with_pymupdf


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".doc",
    ".xlsx",
    ".xls",
    ".csv",
    ".pptx",
    ".ppt",
    ".txt",
}


def is_supported_document(path: Path) -> bool:
    return path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS


def _page(
    doc_id: str,
    doc_type: str,
    language: str,
    source_file: Path,
    page_no: int,
    title: str,
    content: str,
) -> Page:
    body = content.strip() or "(empty page)"
    return Page(
        page_id=f"{doc_id}_p{page_no}",
        doc_id=doc_id,
        doc_type=doc_type,
        language=language,
        content=f"{title}\n{body}".strip(),
        page_no=page_no,
        source_file=str(source_file),
        metadata={"source_ext": source_file.suffix.lower(), "title": title},
    )


def _chunks(lines: Iterable[str], max_chars: int) -> List[str]:
    chunks: List[str] = []
    current: List[str] = []
    current_len = 0
    for raw in lines:
        line = raw.strip()
        if not line:
            continue
        if current and current_len + len(line) + 1 > max_chars:
            chunks.append("\n".join(current))
            current = []
            current_len = 0
        current.append(line)
        current_len += len(line) + 1
    if current:
        chunks.append("\n".join(current))
    return chunks


def _ingest_docx(path: Path, doc_id: str, doc_type: str, language: str, max_chars: int) -> List[Page]:
    try:
        from docx import Document  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("python-docx is required. Install 'python-docx' to use DOCX ingest.") from exc

    doc = Document(str(path))
    lines: List[str] = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            lines.append(paragraph.text)
    for table_idx, table in enumerate(doc.tables, start=1):
        lines.append(f"[table {table_idx}]")
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))

    chunks = _chunks(lines, max_chars) or ["(empty docx)"]
    return [
        _page(doc_id, doc_type, language, path, i, f"DOCX chunk {i}", chunk)
        for i, chunk in enumerate(chunks, start=1)
    ]


def _cell_text(value) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _ingest_xlsx(path: Path, doc_id: str, doc_type: str, language: str, max_rows: int) -> List[Page]:
    try:
        from openpyxl import load_workbook  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("openpyxl is required. Install 'openpyxl' to use XLSX ingest.") from exc

    # openpyxl 对 xlsx 内「数据验证」等扩展会发 UserWarning，与单元格文本抽取无关，避免刷屏
    with warnings.catch_warnings():
        warnings.filterwarnings(
            "ignore",
            message=".*Data Validation extension.*",
            category=UserWarning,
        )
        workbook = load_workbook(str(path), read_only=True, data_only=True)
    pages: List[Page] = []
    for sheet_idx, sheet in enumerate(workbook.worksheets, start=1):
        lines: List[str] = [f"Sheet: {sheet.title}"]
        for row_no, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if row_no > max_rows:
                lines.append(f"... truncated after {max_rows} rows")
                break
            values = [_cell_text(value) for value in row]
            while values and not values[-1]:
                values.pop()
            if values:
                lines.append(" | ".join(values))
        pages.append(
            _page(
                doc_id,
                doc_type,
                language,
                path,
                sheet_idx,
                f"XLSX sheet {sheet_idx}: {sheet.title}",
                "\n".join(lines),
            )
        )
    workbook.close()
    return pages or [_page(doc_id, doc_type, language, path, 1, "XLSX sheet 1", "(empty workbook)")]


def _ingest_csv(path: Path, doc_id: str, doc_type: str, language: str, max_rows: int) -> List[Page]:
    lines: List[str] = []
    with path.open("r", encoding="utf-8-sig", errors="ignore", newline="") as fp:
        reader = csv.reader(fp)
        for row_no, row in enumerate(reader, start=1):
            if row_no > max_rows:
                lines.append(f"... truncated after {max_rows} rows")
                break
            if any(cell.strip() for cell in row):
                lines.append(" | ".join(cell.strip() for cell in row))
    return [_page(doc_id, doc_type, language, path, 1, "CSV table", "\n".join(lines))]


def _ingest_pptx(path: Path, doc_id: str, doc_type: str, language: str) -> List[Page]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("python-pptx is required. Install 'python-pptx' to use PPTX ingest.") from exc

    presentation = Presentation(str(path))
    pages: List[Page] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        lines: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
        pages.append(
            _page(
                doc_id,
                doc_type,
                language,
                path,
                slide_idx,
                f"PPT slide {slide_idx}",
                "\n".join(lines),
            )
        )
    return pages or [_page(doc_id, doc_type, language, path, 1, "PPT slide 1", "(empty pptx)")]


def _ingest_txt(path: Path, doc_id: str, doc_type: str, language: str, max_chars: int) -> List[Page]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    chunks = _chunks(text.splitlines(), max_chars) or ["(empty text file)"]
    return [
        _page(doc_id, doc_type, language, path, i, f"TXT chunk {i}", chunk)
        for i, chunk in enumerate(chunks, start=1)
    ]


def _libreoffice_binary() -> Optional[str]:
    for candidate in ("soffice", "libreoffice"):
        resolved = shutil.which(candidate)
        if resolved:
            return resolved
    mac_app = Path("/Applications/LibreOffice.app/Contents/MacOS/soffice")
    if mac_app.exists():
        return str(mac_app)
    return None


def _convert_to_pdf_with_libreoffice(path: Path, output_dir: Path) -> Path:
    binary = _libreoffice_binary()
    if not binary:
        raise RuntimeError(
            "LibreOffice is required to ingest legacy .doc/.xls/.ppt files. "
            "Install LibreOffice or save the file as .docx/.xlsx/.pptx/.pdf."
        )
    result = subprocess.run(
        [binary, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(path)],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed for {path}: {result.stderr or result.stdout}")
    converted = output_dir / f"{path.stem}.pdf"
    if not converted.exists():
        candidates = list(output_dir.glob("*.pdf"))
        if not candidates:
            raise RuntimeError(f"LibreOffice conversion produced no PDF for {path}")
        converted = candidates[0]
    return converted


def _ingest_via_libreoffice_pdf(
    path: Path,
    doc_id: str,
    doc_type: str,
    language: str,
    image_output_dir: Optional[str],
    dpi: int,
) -> List[Page]:
    with tempfile.TemporaryDirectory(prefix="rag_office_") as tmp:
        pdf_path = _convert_to_pdf_with_libreoffice(path, Path(tmp))
        pages = ingest_pdf_with_pymupdf(
            pdf_path=str(pdf_path),
            doc_id=doc_id,
            doc_type=doc_type,
            language=language,
            image_output_dir=image_output_dir,
            dpi=dpi,
        )
    for page in pages:
        page.source_file = str(path)
        page.metadata["source_ext"] = path.suffix.lower()
        page.metadata["converted_from"] = path.suffix.lower()
    return pages


def ingest_document(
    file_path: str,
    doc_id: str,
    doc_type: str = "manual",
    language: str = "zh",
    image_output_dir: Optional[str] = None,
    dpi: int = 200,
    max_text_chars: int = 4000,
    max_sheet_rows: int = 300,
) -> List[Page]:
    """按扩展名把不同格式统一转换成 Page 列表。"""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"Document not found: {file_path}")

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return ingest_pdf_with_pymupdf(
            pdf_path=str(path),
            doc_id=doc_id,
            doc_type=doc_type,
            language=language,
            image_output_dir=image_output_dir,
            dpi=dpi,
        )
    if suffix == ".docx":
        return _ingest_docx(path, doc_id, doc_type, language, max_text_chars)
    if suffix == ".xlsx":
        return _ingest_xlsx(path, doc_id, doc_type, language, max_sheet_rows)
    if suffix == ".csv":
        return _ingest_csv(path, doc_id, doc_type, language, max_sheet_rows)
    if suffix == ".pptx":
        return _ingest_pptx(path, doc_id, doc_type, language)
    if suffix == ".txt":
        return _ingest_txt(path, doc_id, doc_type, language, max_text_chars)
    if suffix in {".doc", ".xls", ".ppt"}:
        return _ingest_via_libreoffice_pdf(path, doc_id, doc_type, language, image_output_dir, dpi)

    raise ValueError(f"Unsupported document type: {path.suffix}")
